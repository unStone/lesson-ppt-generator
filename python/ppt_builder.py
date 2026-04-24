"""
Enhanced PPT Builder: Generate .pptx from SlidePlan array.
Supports multiple layouts, math formulas, color themes, and animation annotations.
"""
from io import BytesIO
from typing import List
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

from models.schemas import SlidePlan, SlideLayout, ColorScheme, ImageAsset


class PPTBuilder:
    """Build professional PPTX from slide plans."""

    # Layout template names
    LAYOUTS = {
        "title_content": (0, 1),      # title + content
        "title_only": (5, 0),          # title only (used for cover)
        "two_column": (0, 1),          # title + content (we split content)
        "image_text": (0, 1),          # title + content
        "formula": (0, 1),             # title + content
        "chart": (0, 1),               # title + content
        "mindmap": (0, 1),             # title + content
    }

    def __init__(self, color_scheme: ColorScheme = None):
        self.prs = Presentation()
        self.color_scheme = color_scheme or ColorScheme(
            primary="#2563EB", secondary="#DBEAFE",
            background="#FFFFFF", text="#1F2937"
        )
        self._set_slide_size()

    def _set_slide_size(self):
        """Set slide size to widescreen 16:9."""
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def _hex_to_rgb(self, hex_color: str) -> tuple:
        """Convert hex color to RGB tuple."""
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    def _set_shape_color(self, shape, hex_color: str):
        """Set solid fill color for a shape."""
        r, g, b = self._hex_to_rgb(hex_color)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RgbColor(r, g, b)

    def _add_title_slide(self, title: str, subtitle_lines: List[str]):
        """Add a title/cover slide."""
        slide_layout = self.prs.slide_layouts[6]  # blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Background shape (full slide)
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        self._set_shape_color(bg, self.color_scheme.secondary)
        bg.fill.background()  # actually we want solid
        self._set_shape_color(bg, self.color_scheme.secondary)
        # Send to back
        spTree = slide.shapes._spTree
        sp = bg._element
        spTree.remove(sp)
        spTree.insert(2, sp)

        # Accent bar at top
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(0.15)
        )
        self._set_shape_color(accent, self.color_scheme.primary)

        # Title text box
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RgbColor(*self._hex_to_rgb(self.color_scheme.text))
        p.alignment = PP_ALIGN.CENTER

        # Subtitle lines
        if subtitle_lines:
            sub_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2), Inches(12.3), Inches(1.5)
            )
            tf = sub_box.text_frame
            tf.word_wrap = True
            for i, line in enumerate(subtitle_lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(24)
                p.font.color.rgb = RgbColor(*self._hex_to_rgb(self.color_scheme.text))
                p.alignment = PP_ALIGN.CENTER

        # Animation annotation (stored in notes)
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
        else:
            notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = "ANIMATION: fade title, fly_in subtitle lines"

        return slide

    def _add_content_slide(self, plan: SlidePlan):
        """Add a content slide based on layout template."""
        layout_name = plan.layout.layout_template if plan.layout else "title_content"
        slide_layout = self.prs.slide_layouts[0]  # title + content
        slide = self.prs.slides.add_slide(slide_layout)

        # Background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        self._set_shape_color(bg, self.color_scheme.background)
        spTree = slide.shapes._spTree
        sp = bg._element
        spTree.remove(sp)
        spTree.insert(2, sp)

        # Title
        if slide.shapes.title:
            slide.shapes.title.text = plan.title or ""
            for paragraph in slide.shapes.title.text_frame.paragraphs:
                paragraph.font.size = Pt(36)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RgbColor(*self._hex_to_rgb(self.color_scheme.text))

        # Content area
        content_shape = None
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 2:  # body
                content_shape = shape
                break

        if content_shape and plan.content:
            tf = content_shape.text_frame
            tf.word_wrap = True
            tf.clear()

            # Add body text
            for i, line in enumerate(plan.content.body_text or []):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(24)
                p.font.color.rgb = RgbColor(*self._hex_to_rgb(self.color_scheme.text))
                p.space_after = Pt(12)
                # Bullet points for lines starting with bullet-like markers
                if line.startswith("\u2022") or line.startswith("-"):
                    p.level = 0
                elif line.startswith("  ") or line.startswith("\t"):
                    p.level = 1

            # Add formula annotation if present
            if plan.content.formulas:
                p = tf.add_paragraph()
                p.text = ""
                p = tf.add_paragraph()
                p.text = "[FORMULA: " + " | ".join(plan.content.formulas) + "]"
                p.font.size = Pt(18)
                p.font.color.rgb = RgbColor(*self._hex_to_rgb(self.color_scheme.primary))
                p.font.italic = True

            # Add examples if present
            if plan.content.examples:
                for ex in plan.content.examples:
                    p = tf.add_paragraph()
                    p.text = ""
                    p = tf.add_paragraph()
                    p.text = f"[EXAMPLE] {ex.problem}"
                    p.font.size = Pt(20)
                    p.font.bold = True
                    p.font.color.rgb = RgbColor(*self._hex_to_rgb(self.color_scheme.primary))
                    for step in ex.solution_steps:
                        p = tf.add_paragraph()
                        p.text = f"  \u2192 {step}"
                        p.font.size = Pt(18)
                        p.font.color.rgb = RgbColor(*self._hex_to_rgb(self.color_scheme.text))

            # Add diagram annotation
            if plan.content.diagrams_needed:
                p = tf.add_paragraph()
                p.text = ""
                p = tf.add_paragraph()
                p.text = f"[DIAGRAM: {plan.content.diagrams_needed}]"
                p.font.size = Pt(16)
                p.font.color.rgb = RgbColor(255, 0, 0)
                p.font.italic = True

        # Add image placeholder if needed
        if plan.visual and plan.visual.images:
            for img in plan.visual.images:
                # Add a placeholder rectangle with label
                placeholder = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(9), Inches(2), Inches(3.5), Inches(4)
                )
                self._set_shape_color(placeholder, self.color_scheme.secondary)
                placeholder.line.color.rgb = RgbColor(*self._hex_to_rgb(self.color_scheme.primary))
                placeholder.line.width = Pt(2)

                # Add text inside placeholder
                ph_tf = placeholder.text_frame
                ph_tf.word_wrap = True
                p = ph_tf.paragraphs[0]
                p.text = f"[IMAGE: {img.description}]"
                p.font.size = Pt(14)
                p.font.color.rgb = RgbColor(*self._hex_to_rgb(self.color_scheme.primary))
                p.alignment = PP_ALIGN.CENTER

        # Add animation notes
        if plan.animation:
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            anim_notes = []
            for anim in plan.animation.entrance_animations:
                anim_notes.append(f"Element {anim.element_index}: {anim.animation_type} ({anim.trigger})")
            if plan.animation.transition_type:
                anim_notes.append(f"Transition: {plan.animation.transition_type}")
            notes_text_frame.text = "ANIMATION:\n" + "\n".join(anim_notes)

        return slide

    def _add_mindmap_slide(self, plan: SlidePlan):
        """Add a mindmap-style summary slide."""
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)

        # Background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        self._set_shape_color(bg, self.color_scheme.background)
        spTree = slide.shapes._spTree
        sp = bg._element
        spTree.remove(sp)
        spTree.insert(2, sp)

        # Title
        if slide.shapes.title:
            slide.shapes.title.text = plan.title or "课堂小结"
            for paragraph in slide.shapes.title.text_frame.paragraphs:
                paragraph.font.size = Pt(36)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RgbColor(*self._hex_to_rgb(self.color_scheme.text))

        # Central node (topic)
        center = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(5.5), Inches(3), Inches(2.3), Inches(1.5)
        )
        self._set_shape_color(center, self.color_scheme.primary)
        center_tf = center.text_frame
        center_tf.word_wrap = True
        p = center_tf.paragraphs[0]
        p.text = plan.content.title or "主题"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RgbColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Branch nodes from body text
        if plan.content and plan.content.body_text:
            positions = [
                (Inches(2), Inches(1.5)), (Inches(9.5), Inches(1.5)),
                (Inches(2), Inches(5)), (Inches(9.5), Inches(5)),
                (Inches(0.5), Inches(3.25)), (Inches(11.5), Inches(3.25)),
            ]
            for i, line in enumerate(plan.content.body_text[:6]):
                if not line.strip():
                    continue
                x, y = positions[i % len(positions)]
                branch = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.8), Inches(0.8)
                )
                self._set_shape_color(branch, self.color_scheme.secondary)
                branch_tf = branch.text_frame
                branch_tf.word_wrap = True
                p = branch_tf.paragraphs[0]
                p.text = line[:30]  # Truncate long lines
                p.font.size = Pt(16)
                p.font.color.rgb = RgbColor(*self._hex_to_rgb(self.color_scheme.text))
                p.alignment = PP_ALIGN.CENTER

        return slide

    def build(self, slides: List[SlidePlan]) -> bytes:
        """Build complete PPTX and return as bytes."""
        for plan in slides:
            if plan.section == "封面" or plan.layout and plan.layout.layout_template == "title_only":
                self._add_title_slide(
                    plan.title,
                    plan.content.body_text if plan.content else []
                )
            elif plan.layout and plan.layout.layout_template == "mindmap":
                self._add_mindmap_slide(plan)
            else:
                self._add_content_slide(plan)

        # Save to bytes
        output = BytesIO()
        self.prs.save(output)
        output.seek(0)
        return output.getvalue()
